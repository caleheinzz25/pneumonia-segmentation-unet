#!/usr/bin/env python3
"""
=============================================================================
NATIVE POWERPOINT (PPTX) GENERATOR — GAYA UNIVERSITAS GUNADARMA
Generates 100% editable 15-slide presentation.pptx matching PDF 1:1 Clean Beamer Style
=============================================================================
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ─── PATHS ───────────────────────────────────────────────────────────────────
SCRIPT_DIR  = os.path.dirname(os.path.abspath(__file__))
IMG_DIR     = os.path.join(os.path.dirname(SCRIPT_DIR), "img")
OUT_FILE    = os.path.join(SCRIPT_DIR, "presentation.pptx")

# ─── COLOUR PALETTE (Gunadarma / Beamer Theme) ───────────────────────────────
C_MAGENTA    = RGBColor(188, 36, 186)   # #BC24BA
C_PURPLE     = RGBColor(147, 51, 234)   # #9333EA
C_HEADER_R   = RGBColor(160, 30, 160)   # Lighter top-right header box
C_LIGHT_PINK = RGBColor(245, 210, 245)  # #F5D2F5
C_WHITE      = RGBColor(255, 255, 255)
C_BLACK      = RGBColor(20, 20, 20)
C_GREY       = RGBColor(80, 80, 80)
C_ALERT_BG   = RGBColor(253, 238, 245)
C_ALERT_TXT  = RGBColor(219, 39, 119)   # Magenta-red for alert

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_FAMILY = "Arial"


def img(name):
    p = os.path.join(IMG_DIR, name)
    return p if os.path.exists(p) else None


def add_slide_header(slide, title, subtitle=""):
    # Main Magenta Header Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_MAGENTA
    bar.line.fill.background()

    # Right Header Box (Lighter Purple/Magenta Block)
    r_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.0), 0, Inches(2.333), Inches(1.15))
    r_box.fill.solid()
    r_box.fill.fore_color.rgb = C_HEADER_R
    r_box.line.fill.background()

    # Header Textbox Left
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(10.4), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0

    p_t = tf.paragraphs[0]
    p_t.text = title
    p_t.font.name = FONT_FAMILY
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = C_WHITE

    if subtitle:
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = FONT_FAMILY
        p_sub.font.size = Pt(11)
        p_sub.font.italic = True
        p_sub.font.color.rgb = C_LIGHT_PINK

    # Header Right Text
    tb_r = slide.shapes.add_textbox(Inches(11.1), Inches(0.20), Inches(2.1), Inches(0.8))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_top = tf_r.margin_bottom = tf_r.margin_left = tf_r.margin_right = 0
    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "Gunadarma"
    p_r1.alignment = PP_ALIGN.RIGHT
    p_r1.font.name = FONT_FAMILY
    p_r1.font.size = Pt(11)
    p_r1.font.bold = True
    p_r1.font.color.rgb = C_WHITE

    p_r2 = tf_r.add_paragraph()
    p_r2.text = "UG University"
    p_r2.alignment = PP_ALIGN.RIGHT
    p_r2.font.name = FONT_FAMILY
    p_r2.font.size = Pt(9.5)
    p_r2.font.color.rgb = C_LIGHT_PINK


def add_slide_footer(slide, slide_num):
    # Bottom Footer Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.95), SLIDE_W, Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_MAGENTA
    bar.line.fill.background()

    # Left text block
    tb_l = slide.shapes.add_textbox(Inches(0.3), Inches(6.98), Inches(3.5), Inches(0.5))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_top = tf_l.margin_bottom = tf_l.margin_left = tf_l.margin_right = 0
    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "visit our website"
    p_l1.font.name = FONT_FAMILY
    p_l1.font.size = Pt(8)
    p_l1.font.bold = True
    p_l1.font.color.rgb = C_WHITE

    p_l2 = tf_l.add_paragraph()
    p_l2.text = "www.gunadarma.ac.id"
    p_l2.font.name = FONT_FAMILY
    p_l2.font.size = Pt(8)
    p_l2.font.color.rgb = C_LIGHT_PINK

    # Center-Right text block
    tb_c = slide.shapes.add_textbox(Inches(4.2), Inches(6.98), Inches(6.8), Inches(0.5))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_top = tf_c.margin_bottom = tf_c.margin_left = tf_c.margin_right = 0
    p_c1 = tf_c.paragraphs[0]
    p_c1.text = "More Information — GUNADARMA UNIVERSITY"
    p_c1.alignment = PP_ALIGN.RIGHT
    p_c1.font.name = FONT_FAMILY
    p_c1.font.size = Pt(8)
    p_c1.font.bold = True
    p_c1.font.color.rgb = C_WHITE

    p_c2 = tf_c.add_paragraph()
    p_c2.text = "Jl. Margonda Raya 100, Depok · Telp. (+62-21) 7888 1112"
    p_c2.alignment = PP_ALIGN.RIGHT
    p_c2.font.name = FONT_FAMILY
    p_c2.font.size = Pt(8)
    p_c2.font.color.rgb = C_LIGHT_PINK

    # Right logo
    logo_p = img("logo_gunadarma.png")
    if logo_p:
        slide.shapes.add_picture(logo_p, Inches(11.2), Inches(6.96), height=Inches(0.50))

    # Page number at bottom left
    tb_num = slide.shapes.add_textbox(Inches(0.3), Inches(6.75), Inches(0.8), Inches(0.3))
    p_num = tb_num.text_frame.paragraphs[0]
    p_num.text = f"{slide_num}"
    p_num.font.name = FONT_FAMILY
    p_num.font.size = Pt(9)
    p_num.font.color.rgb = C_GREY


def add_clean_block(slide, left, top, width, height, title, items, title_color=C_MAGENTA):
    """Clean Beamer-style block with title and bullet items, NO ugly outer box borders."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0

    p_t = tf.paragraphs[0]
    p_t.text = title
    p_t.font.name = FONT_FAMILY
    p_t.font.size = Pt(13)
    p_t.font.bold = True
    p_t.font.color.rgb = title_color
    p_t.space_after = Pt(6)

    for item in items:
        p = tf.add_paragraph()
        p.font.name = FONT_FAMILY
        p.font.size = Pt(10.5)
        p.space_after = Pt(4)

        if ":" in item and not item.startswith("http"):
            parts = item.split(":", 1)
            # Add bullet prefix
            p.text = parts[0] + ":"
            p.font.bold = True
            p.font.color.rgb = C_BLACK
            
            # Add remaining text non-bold
            run = p.add_run()
            run.text = parts[1]
            run.font.bold = False
            run.font.color.rgb = C_BLACK
        else:
            p.text = item
            p.font.bold = False
            p.font.color.rgb = C_BLACK


def add_alert_block(slide, left, top, width, height, title, text_content):
    """Clean alert block with light pink background box matching Beamer AlertBlock."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = C_ALERT_BG
    box.line.fill.background()

    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0

    p_t = tf.paragraphs[0]
    p_t.text = title
    p_t.font.name = FONT_FAMILY
    p_t.font.size = Pt(12)
    p_t.font.bold = True
    p_t.font.color.rgb = C_ALERT_TXT
    p_t.space_after = Pt(4)

    p_b = tf.add_paragraph()
    p_b.text = text_content
    p_b.font.name = FONT_FAMILY
    p_b.font.size = Pt(10)
    p_b.font.color.rgb = C_BLACK


def create_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    # =========================================================================
    # SLIDE 1: JUDUL & IDENTITAS (COVER SLIDE)
    # =========================================================================
    sl1 = prs.slides.add_slide(blank_layout)

    # Top-left Logo
    logo_p = img("logo_gunadarma.png")
    if logo_p:
        sl1.shapes.add_picture(logo_p, Inches(0.4), Inches(0.2), height=Inches(1.3))

    # Middle Split Banner Block
    banner_left = sl1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.2), Inches(6.66), Inches(2.8)
    )
    banner_left.fill.solid()
    banner_left.fill.fore_color.rgb = C_MAGENTA
    banner_left.line.fill.background()

    tb_c = sl1.shapes.add_textbox(Inches(0.4), Inches(2.35), Inches(6.0), Inches(2.5))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    p_ct = tf_c.paragraphs[0]
    p_ct.text = "DETEKSI AREA INFEKSI PNEUMONIA PADA CITRA RONTGEN DADA MENGGUNAKAN ARSITEKTUR U-NET BERBASIS CONVOLUTIONAL NEURAL NETWORK (CNN)"
    p_ct.font.name = FONT_FAMILY
    p_ct.font.size = Pt(18)
    p_ct.font.bold = True
    p_ct.font.color.rgb = C_WHITE

    cxr_p = img("cxr_normal_vs_pneumonia.png")
    if cxr_p:
        sl1.shapes.add_picture(cxr_p, Inches(6.66), Inches(2.2), Inches(6.67), Inches(2.8))

    # Lower Author Info
    tb_au = sl1.shapes.add_textbox(Inches(0.4), Inches(5.2), Inches(7.5), Inches(1.8))
    tf_au = tb_au.text_frame
    tf_au.word_wrap = True

    p_a1 = tf_au.paragraphs[0]
    p_a1.text = "Nama : Muhammad Hisyam"
    p_a1.font.name = FONT_FAMILY
    p_a1.font.size = Pt(15)
    p_a1.font.bold = True

    p_a2 = tf_au.add_paragraph()
    p_a2.text = "NPM : 51422075"
    p_a2.font.name = FONT_FAMILY
    p_a2.font.size = Pt(15)
    p_a2.font.bold = True

    p_a3 = tf_au.add_paragraph()
    p_a3.text = "Pembimbing : Dr. Tavipia Rumambi, SKom., MMSI."
    p_a3.font.name = FONT_FAMILY
    p_a3.font.size = Pt(15)
    p_a3.font.bold = True

    # Bottom Right Logo & Text
    if logo_p:
        sl1.shapes.add_picture(logo_p, Inches(11.8), Inches(5.8), height=Inches(1.1))
    tb_sk = sl1.shapes.add_textbox(Inches(7.5), Inches(6.3), Inches(4.2), Inches(0.5))
    p_sk = tb_sk.text_frame.paragraphs[0]
    p_sk.text = "SKRIPSI UNIVERSITAS GUNADARMA"
    p_sk.alignment = PP_ALIGN.RIGHT
    p_sk.font.name = FONT_FAMILY
    p_sk.font.size = Pt(13)
    p_sk.font.bold = True
    p_sk.font.color.rgb = C_MAGENTA

    # =========================================================================
    # SLIDE 2: LATAR BELAKANG MASALAH
    # =========================================================================
    sl2 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl2, "LATAR BELAKANG MASALAH", "Tantangan Diagnostik Pneumonia & Urgensi Sistem Deteksi Otomatis")
    add_slide_footer(sl2, 2)

    add_clean_block(sl2, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.6),
                    "Beban Klinis & Keterbatasan Diagnosis Manual",
                    ["• Mortalitas Tinggi: Pneumonia merenggut ~2,5 juta jiwa global tahun 2019 (WHO), dengan korban terbesar anak di bawah 5 tahun & lansia.",
                     "• Diagnosis via CXR: Foto rontgen dada dipilih karena murah & dosis radiasi rendah dibanding CT Scan, namun rentan salah baca — terutama di fasilitas minim tenaga ahli.",
                     "• Beban Radiolog: Kelebihan beban kerja radiolog memperlambat diagnosis, sementara pola lesi pneumonia sering mirip edema atau efusi pleura."])

    add_clean_block(sl2, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3.6),
                    "Solusi: Deep Learning & Explainable AI",
                    ["• CNN U-Net: Terbukti mampu segmentasi citra medis piksel-demi-piksel dengan akurasi tinggi.",
                     "• sCSE Attention: Mekanisme atensi memfokuskan model hanya pada area lesi relevan, mengabaikan bagian tidak penting.",
                     "• EfficientNet-B3 (Transfer Learning): Mengatasi keterbatasan data medis — mempercepat pelatihan & meningkatkan performa.",
                     "• Grad-CAM & Gradio Web: Menjamin transparansi keputusan model & kemudahan akses bagi tenaga medis."])

    add_alert_block(sl2, Inches(0.6), Inches(5.2), Inches(12.0), Inches(1.3),
                    "Sistem yang Dibangun",
                    "Penelitian ini membangun sistem CDSS (Clinical Decision Support System) berbasis U-Net + sCSE + EfficientNet-B3 mengikuti kerangka CRISP-DM, dilengkapi Grad-CAM dan antarmuka web Gradio yang dapat diakses via internet melalui Cloudflare Tunnel.")

    # =========================================================================
    # SLIDE 3: URGENSI SEGMENTASI MEDIS OTOMATIS
    # =========================================================================
    sl3 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl3, "URGENSI SEGMENTASI MEDIS OTOMATIS", "Komparasi Paradigma: Klasifikasi vs. Deteksi Objek vs. Segmentasi Semantik (U-Net)")
    add_slide_footer(sl3, 3)

    add_clean_block(sl3, Inches(0.6), Inches(1.4), Inches(3.8), Inches(5.2),
                    "1. Klasifikasi Biner\nResNet / VGG",
                    ["• Output: Label tunggal (Pneumonia vs Normal).",
                     "• Keterbatasan: Tidak memberikan informasi lokasi, batas, maupun luas infeksi.",
                     "• Klinis: Kurang informatif untuk pemantauan progresivitas."])

    add_clean_block(sl3, Inches(4.7), Inches(1.4), Inches(3.8), Inches(5.2),
                    "2. Deteksi Objek\nYOLO / Faster R-CNN",
                    ["• Output: Bounding Box persegi (koordinat x, y, w, h).",
                     "• Keterbatasan: Lesi ireguler. Box menyertakan banyak jaringan sehat.",
                     "• Klinis: Kalkulasi rasio keparahan menjadi bias & tidak akurat."])

    add_clean_block(sl3, Inches(8.8), Inches(1.4), Inches(3.8), Inches(5.2),
                    "3. Segmentasi U-Net\nU-Net + sCSE (Solusi)",
                    ["• Output: Masker biner piksel-demi-piksel (pixel-wise).",
                     "• Keunggulan: Memetakan kontur lesi ireguler secara terperinci.",
                     "• Klinis: Memungkinkan kalkulasi Severity Score (%)."],
                    title_color=C_PURPLE)

    # =========================================================================
    # SLIDE 4: RUMUSAN MASALAH & TUJUAN
    # =========================================================================
    sl4 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl4, "RUMUSAN MASALAH & TUJUAN PENELITIAN", "Fokus Kajian Ilmiah & Target Pencapaian Penelitian")
    add_slide_footer(sl4, 4)

    add_clean_block(sl4, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.2),
                    "Rumusan Masalah Penelitian",
                    ["1. Bagaimana membangun dan melatih model deteksi pneumonia menggunakan arsitektur U-Net dengan atensi sCSE dan encoder EfficientNet-B3 pada dataset RSNA Pneumonia Detection Challenge?",
                     "2. Bagaimana pengaruh penggunaan mekanisme atensi sCSE dan pra-pemrosesan segmentasi paru-paru otomatis (PSPNet) terhadap performa model dalam mendeteksi area infeksi pneumonia?",
                     "3. Bagaimana membangun antarmuka aplikasi web yang dapat menampilkan hasil deteksi beserta visualisasi Grad-CAM agar mudah dipahami oleh pengguna?"])

    add_clean_block(sl4, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.2),
                    "Tujuan Penelitian",
                    ["1. Membangun dan melatih model deteksi pneumonia berbasis U-Net dengan atensi sCSE dan encoder EfficientNet-B3 pada dataset RSNA Pneumonia Detection Challenge.",
                     "2. Mengevaluasi performa model menggunakan metrik segmentasi (Dice Coefficient, Precision, Recall, Specificity) serta menganalisis pengaruh sCSE dan PSPNet terhadap hasil deteksi.",
                     "3. Membangun aplikasi web berbasis Gradio yang mengintegrasikan model deteksi dengan visualisasi Grad-CAM agar hasil mudah dipahami pengguna."],
                    title_color=C_PURPLE)

    # =========================================================================
    # SLIDE 5: KERANGKA KERJA CRISP-DM
    # =========================================================================
    sl5 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl5, "KERANGKA KERJA CRISP-DM", "Siklus Metodologi Pengembangan Sistem (6 Fase Utama)")
    add_slide_footer(sl5, 5)

    crisp_p = img("crisp_dm.png")
    if crisp_p:
        sl5.shapes.add_picture(crisp_p, Inches(0.6), Inches(1.4), Inches(4.5), Inches(5.2))

    add_clean_block(sl5, Inches(5.4), Inches(1.4), Inches(7.3), Inches(5.2),
                    "6 Fase CRISP-DM pada Penelitian Ini",
                    ["1. Business Understanding: Kebutuhan CDSS screening awal penunjang keputusan dokter.",
                     "2. Data Understanding: Dataset RSNA Pneumonia (30.227 CXR DICOM), analisis class imbalance.",
                     "3. Data Preparation: Konversi DICOM ke PNG 8-bit, CLAHE enhancement, PSPNet Dual Lung Masking, augmentasi data (8 teknik).",
                     "4. Modeling: U-Net + EfficientNet-B3 + sCSE + Unified Focal Loss + pelatihan 2-fase (AdamW, AMP).",
                     "5. Evaluation: Evaluasi kuantitatif (Dice 0.6234, Recall 0.7082) & validasi visual Grad-CAM.",
                     "6. Deployment: Gradio Web App 4-panel & Cloudflare Tunnel (grad.mhisyam.com)."])

    # =========================================================================
    # SLIDE 6: FLOWCHART PENELITIAN
    # =========================================================================
    sl6 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl6, "FLOWCHART PENELITIAN", "Alur Tahapan Eksperimen & Pemrosesan Data dari Input hingga Web Deployment")
    add_slide_footer(sl6, 6)

    flow_p = img("flowchart_penelitian.png")
    if flow_p:
        sl6.shapes.add_picture(flow_p, Inches(0.8), Inches(1.4), Inches(3.6), Inches(5.2))

    add_clean_block(sl6, Inches(4.8), Inches(1.4), Inches(7.8), Inches(5.2),
                    "Rincian Tahapan Alur Penelitian",
                    ["• Input Data: Citra Rontgen Dada DICOM 16-bit dari RSNA Pneumonia Detection Challenge.",
                     "• Preprocessing: Normalisasi Min-Max, kontras CLAHE, & segmentasi ROI paru PSPNet.",
                     "• Augmentasi Data: 8 teknik Albumentations (Rotation, Shift, Brightness, Noise, Dropout).",
                     "• Training U-Net: Optimizer AdamW, AMP FP16, Unified Focal Loss (α=0.5, γ=2.0).",
                     "• Evaluasi Model: Pengujian pada 3.543 citra test set & Confusion Matrix piksel.",
                     "• XAI & Web App: Peta gradien Grad-CAM, kalkulasi Severity Tier, & Gradio Web deployment."])

    # =========================================================================
    # SLIDE 7: DIAGRAM PERANCANGAN SISTEM (UML)
    # =========================================================================
    sl7 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl7, "DIAGRAM PERANCANGAN SISTEM (UML)", "Model Interaksi Pengguna (Use Case) & Alur Logika Inferensi (Activity Diagram)")
    add_slide_footer(sl7, 7)

    uc_p = img("use_case_diagram.png")
    if uc_p:
        sl7.shapes.add_picture(uc_p, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.2))

    act_p = img("activity_diagram.png")
    if act_p:
        sl7.shapes.add_picture(act_p, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3.2))

    add_clean_block(sl7, Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.0),
                    "Use Case Diagram Aplikasi Web",
                    ["• 6 Fungsi Utama: Unggah Citra (DICOM/PNG/JPG), Atur Threshold Slider, Analisis File U-Net, Lihat Overlay & Grad-CAM, Lihat Laporan Severity, dan Reset."])

    add_clean_block(sl7, Inches(6.8), Inches(4.7), Inches(5.8), Inches(2.0),
                    "Activity Diagram Inferensi Medis",
                    ["• Alur Logika: Input Rontgen ➔ Preprocessing CLAHE & PSPNet ➔ Inferensi U-Net + sCSE ➔ Generasi Heatmap Grad-CAM ➔ Output 4-Panel & Severity Report."])

    # =========================================================================
    # SLIDE 8: DATASET & PREPROCESSING
    # =========================================================================
    sl8 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl8, "DATASET & PREPROCESSING PIPELINE", "RSNA Pneumonia Detection Challenge & Dual Lung Masking")
    add_slide_footer(sl8, 8)

    add_clean_block(sl8, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.2),
                    "Dataset RSNA Pneumonia Challenge",
                    ["• Total Citra: 26.684 citra latih DICOM (1024 x 1024).",
                     "• Lung Opacity (Pneumonia Positif): 6.012 citra (22,5%).",
                     "• No Lung Opacity / Not Normal: 7.210 citra (27,1%).",
                     "• Normal: 13.462 citra (50,4%).",
                     "• Patient-Level Split: Split 80:20 pada tingkat Patient ID (mencegah data leakage)."])

    add_clean_block(sl8, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.2),
                    "Preprocessing & Dual Lung Masking",
                    ["• 16-bit ke 8-bit: Normalisasi Min-Max [0, 255] & resize 512 x 512.",
                     "• CLAHE: Peningkatan kontras lokal (clip_limit=2.0, tile 8x8).",
                     "• PSPNet Dual Lung Masking: Memotong ROI paru kiri & kanan, mengeliminasi teks DICOM, kabel EKG, & penanda L/R.",
                     "• 8 Augmentasi: Horizontal Flip, Shift-Scale-Rotate, Gamma, Noise, Coarse/Grid Dropout."])

    # =========================================================================
    # SLIDE 8: KONSEP DASAR CONVOLUTIONAL NEURAL NETWORK (CNN)
    # =========================================================================
    sl8 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl8, "KONSEP DASAR CONVOLUTIONAL NEURAL NETWORK (CNN)", "Prinsip Ekstraksi Fitur Citra: Operasi Konvolusi, Fungsi Aktivasi, & Pooling")
    add_slide_footer(sl8, 8)

    add_clean_block(sl8, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.2),
                    "Komponen Utama Arsitektur CNN",
                    ["• Operasi Konvolusi (Y = f(W * X + b)): Perkalian kernel filter terhadap piksel citra untuk mengekstrak fitur spasial (tepi, tekstur, opasitas).",
                     "• Fungsi Aktivasi (ReLU / Swish): Transformasi non-linear (max(0, x)) agar jaring mempelajari pola lesi ireguler yang kompleks.",
                     "• Pooling Layer (Max/Average): Reduksi dimensi spasial (downsampling) untuk efisiensi komputasi & ketahanan variasi posisi.",
                     "• Feature Maps Hirarkis: Ekstraksi dari fitur tingkat rendah (edges) hingga tingkat tinggi (infiltrat pneumonia)."])

    conv_p = img("operasi_konvolusi.png")
    if conv_p:
        sl8.shapes.add_picture(conv_p, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.7))

    add_clean_block(sl8, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.3),
                    "Keunggulan CNN pada Segmentasi Medis",
                    ["• Parameter Sharing: Penggunaan ulang bobot kernel di seluruh piksel citra, meminimalkan risiko overfitting.",
                     "• Spatial Hierarchy: Mempertahankan relasi spasial antar-piksel tetangga pada organ paru.",
                     "• Pretrained Backbone: Menjadi fondasi Feature Extractor pada arsitektur U-Net (EfficientNet-B3)."])

    # =========================================================================
    # SLIDE 9: ARSITEKTUR U-NET & FORMULA (HIGH-RES TIKZ DIAGRAM)
    # =========================================================================
    sl9 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl9, "ARSITEKTUR U-NET & FORMULASI MATEMATIKA", "U-Net + EfficientNet-B3 Backbone + sCSE Attention Gate + Unified Focal Loss")
    add_slide_footer(sl9, 9)

    tikz_img = img("unet_architecture_tikz.png")
    if tikz_img:
        sl9.shapes.add_picture(tikz_img, Inches(0.5), Inches(1.4), Inches(6.2), Inches(5.2))

    add_clean_block(sl9, Inches(6.9), Inches(1.4), Inches(5.8), Inches(2.5),
                    "Formulasi Unified Focal Loss",
                    ["L_UFL = 0.5 * L_Focal + 0.5 * L_FT",
                     "L_Focal = -alpha_t * (1 - p_t)^gamma * log(p_t)",
                     "Parameter Optimal: alpha = 0.5, gamma = 2.0"])

    add_clean_block(sl9, Inches(6.9), Inches(4.1), Inches(5.8), Inches(2.5),
                    "Formulasi Atensi sCSE",
                    ["sCSE(F) = cSE(F) + sSE(F)",
                     "cSE = sigma(MLP(GAP(F))) * F",
                     "sSE = sigma(Conv1x1(F)) * F"])

    # =========================================================================
    # SLIDE 10: HASIL EVALUASI KUANTITATIF & CONFUSION MATRIX
    # =========================================================================
    sl10 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl10, "HASIL EVALUASI KUANTITATIF & CONFUSION MATRIX", "Evaluasi Performa Segmentasi pada Data Validasi (350.224.384 Total Piksel)")
    add_slide_footer(sl10, 10)

    add_clean_block(sl10, Inches(0.6), Inches(1.4), Inches(5.6), Inches(5.2),
                    "Metrik Utama Performa (KPI Target ≥ 0,60)",
                    ["• Dice Coefficient: 0.6234 (62,34%)",
                     "• Precision: 0.6868 (68,68%)",
                     "• Recall (Sensitivitas): 0.7082 (70,82%)",
                     "• Specificity: 0.9781 (97,81%)",
                     "• Validation Loss: 0.2755 (Epoch 46)"])

    rows, cols = 3, 3
    left, top, width, height = Inches(6.5), Inches(1.4), Inches(6.2), Inches(2.3)
    table_shape = sl10.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    cm_headers = ["", "Pred. Normal", "Pred. Pneumonia"]
    cm_data = [
        ["Act. Normal", "TN = 312.284.691 (97,81%)", "FP = 7.198.145 (2,19%)"],
        ["Act. Pneumonia", "FN = 9.995.562 (29,18%)", "TP = 20.745.986 (70,82%)"]
    ]

    for c_idx, h in enumerate(cm_headers):
        cell = table.cell(0, c_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_MAGENTA
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_FAMILY
            p.font.size = Pt(10.5)
            p.font.bold = True
            p.font.color.rgb = C_WHITE

    for r_idx, row_data in enumerate(cm_data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_FAMILY
                p.font.size = Pt(10)
                if c_idx == 0:
                    p.font.bold = True
                p.font.color.rgb = C_BLACK

    add_alert_block(sl10, Inches(6.5), Inches(4.2), Inches(6.2), Inches(2.4),
                    "Signifikansi Klinis",
                    "Recall 70,82% meminimalkan risiko false negative medis, sementara Specificity 97,81% mencegah false alarm berlebih.")

    # =========================================================================
    # SLIDE 11: PERBANDINGAN BASELINE & ABLATION STUDY
    # =========================================================================
    sl11 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl11, "PERBANDINGAN BASELINE & ABLATION STUDY", "Analisis Kontribusi Inkremental Komponen Model (+14,1% Dice Total Gain)")
    add_slide_footer(sl11, 11)

    rows, cols = 6, 5
    left, top, width, height = Inches(0.6), Inches(1.4), Inches(12.0), Inches(2.8)
    table_shape = sl11.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["Konfigurasi Komponen Model", "Dice", "Recall", "Precision", "Δ Dice"]
    data = [
        ["Baseline U-Net (Tanpa Pretrain)", "0.4820", "0.5310", "0.5240", "—"],
        ["+ EfficientNet-B3 Pretrained ImageNet", "0.5410", "0.6070", "0.5730", "+5,9%"],
        ["+ sCSE Attention Gate", "0.5890", "0.6620", "0.6280", "+4,8%"],
        ["+ Unified Focal Loss (α=0.5, γ=2.0)", "0.6080", "0.6940", "0.6610", "+1,9%"],
        ["+ CLAHE & Augmentasi Data (Model Final)", "0.6234", "0.7082", "0.6868", "+1,5%"]
    ]

    for c_idx, h in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_MAGENTA
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_FAMILY
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = C_WHITE

    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_LIGHT_PINK if r_idx == 4 else C_WHITE
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_FAMILY
                p.font.size = Pt(10)
                if r_idx == 4:
                    p.font.bold = True
                p.font.color.rgb = C_BLACK

    add_clean_block(sl11, Inches(0.6), Inches(4.5), Inches(12.0), Inches(2.1),
                    "Key Takeaways Ablation Study",
                    ["• Pemanfaatan bobot ImageNet pada encoder EfficientNet-B3 memberikan lonjakan terbesar (+5,9% Dice).",
                     "• Atensi sCSE memperjelas batas lesi & menyaring fitur relevan pada skip connections (+4,8% Dice).",
                     "• Total peningkatan kumulatif Dice Score mencapai +14,1% dibanding U-Net baseline standar (0.4820 ke 0.6234)."])

    # =========================================================================
    # SLIDE 12: VISUALISASI SEGMENTASI KLINIS
    # =========================================================================
    sl12 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl12, "VISUALISASI SEGMENTASI KLINIS", "Evaluasi 4 Sampel Kasus Medis (4 Kolom: CXR, PSPNet Masking, Overlay, Heatmap)")
    add_slide_footer(sl12, 12)

    vis_p = img("hasil_segmentasi_visual.png")
    if vis_p:
        sl12.shapes.add_picture(vis_p, Inches(0.8), Inches(1.4), Inches(11.7), Inches(4.0))

    col_desc = [
        ("(a) Citra CXR Asli", "Rontgen mentah pasien"),
        ("(b) PSPNet Lung Mask", "Hasil pemotongan ROI organ"),
        ("(c) Overlay Prediction", "Masker U-Net (Arsir Merah)"),
        ("(d) Grad-CAM Heatmap", "Peta aktivasi JET colormap")
    ]

    for idx, (title, desc) in enumerate(col_desc):
        c_left = Inches(0.8 + idx * 2.95)
        tb_col = sl12.shapes.add_textbox(c_left, Inches(5.5), Inches(2.8), Inches(1.1))
        tf_col = tb_col.text_frame
        tf_col.word_wrap = True
        p1 = tf_col.paragraphs[0]
        p1.text = title
        p1.font.name = FONT_FAMILY
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.alignment = PP_ALIGN.CENTER
        p1.font.color.rgb = C_MAGENTA

        p2 = tf_col.add_paragraph()
        p2.text = desc
        p2.font.name = FONT_FAMILY
        p2.font.size = Pt(9)
        p2.alignment = PP_ALIGN.CENTER
        p2.font.color.rgb = C_GREY

    # =========================================================================
    # SLIDE 13: EXPLAINABLE AI (GRAD-CAM & SEVERITY)
    # =========================================================================
    sl13 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl13, "EXPLAINABLE AI (XAI) & SEVERITY SCORING", "Interpretasi Medis Heatmap Grad-CAM & Tingkat Keparahan Infeksi")
    add_slide_footer(sl13, 13)

    rows, cols = 5, 3
    left, top, width, height = Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.2)
    table_shape = sl13.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    jet_headers = ["Warna", "Probabilitas", "Interpretasi Klinis"]
    jet_data = [
        ["Merah", "0,8 - 1,0", "Pusat lesi infeksi / gradien dominan"],
        ["Kuning", "0,6 - 0,8", "Aktivasi tinggi / margin lesi transisi"],
        ["Hijau", "0,3 - 0,6", "Aktivasi sedang / perilesional"],
        ["Biru", "0,0 - 0,3", "Jaringan paru sehat / background"]
    ]
    bg_colors = [RGBColor(254, 226, 226), RGBColor(254, 249, 195), RGBColor(220, 252, 231), RGBColor(224, 242, 254)]

    for c_idx, h in enumerate(jet_headers):
        cell = table.cell(0, c_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_MAGENTA
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_FAMILY
            p.font.size = Pt(10.5)
            p.font.bold = True
            p.font.color.rgb = C_WHITE

    for r_idx, row_data in enumerate(jet_data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_colors[r_idx]
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_FAMILY
                p.font.size = Pt(9.5)
                if c_idx == 0:
                    p.font.bold = True
                p.font.color.rgb = C_BLACK

    add_clean_block(sl13, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.2),
                    "5-Tier Severity Scoring System",
                    ["• Level 0 (Normal): 0% rasio infeksi.",
                     "• Level 1 (Mild): ≤ 5% luas organ paru.",
                     "• Level 2 (Moderate): 5% - 15% luas organ paru.",
                     "• Level 3 (Severe): 15% - 30% luas organ paru.",
                     "• Level 4 (Critical): > 30% luas organ paru.",
                     "• Lateralitas: Dextra, Sinistra, Bilateral."],
                    title_color=C_PURPLE)

    # =========================================================================
    # SLIDE 14: DIAGRAM PERANCANGAN SISTEM (UML)
    # =========================================================================
    sl14 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl14, "DIAGRAM PERANCANGAN SISTEM (UML)", "Model Interaksi Pengguna (Use Case) & Alur Logika Inferensi (Activity Diagram)")
    add_slide_footer(sl14, 14)

    uc_p = img("use_case_diagram.png")
    if uc_p:
        sl14.shapes.add_picture(uc_p, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.2))

    act_p = img("activity_diagram.png")
    if act_p:
        sl14.shapes.add_picture(act_p, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3.2))

    add_clean_block(sl14, Inches(0.6), Inches(4.7), Inches(5.8), Inches(2.0),
                    "Use Case Diagram Aplikasi Web",
                    ["• 6 Fungsi Utama: Unggah Citra (DICOM/PNG/JPG), Atur Threshold Slider, Analisis File U-Net, Lihat Overlay & Grad-CAM, Lihat Laporan Severity, dan Reset."])

    add_clean_block(sl14, Inches(6.8), Inches(4.7), Inches(5.8), Inches(2.0),
                    "Activity Diagram Inferensi Medis",
                    ["• Alur Logika: Input Rontgen ➔ Preprocessing CLAHE & PSPNet ➔ Inferensi U-Net + sCSE ➔ Generasi Heatmap Grad-CAM ➔ Output 4-Panel & Severity Report."])

    # =========================================================================
    # SLIDE 15: TAMPILAN WEBSITE & CLOUDFLARE TUNNEL
    # =========================================================================
    sl15 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl15, "TAMPILAN WEBSITE & CLOUDFLARE TUNNEL DEPLOYMENT", "Antarmuka Web Gradio Interaktif & Akses Publik Terenkripsi HTTPS")
    add_slide_footer(sl15, 15)

    grad_p = img("implementasi_web_gradio.png")
    if grad_p:
        sl15.shapes.add_picture(grad_p, Inches(0.6), Inches(1.4), Inches(5.8), Inches(3.0))

    cf_p = img("cloudflare_tunnel.png")
    if cf_p:
        sl15.shapes.add_picture(cf_p, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3.0))

    add_clean_block(sl15, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.1),
                    "Fitur Gradio Web App",
                    ["• Dual Input (DICOM / PNG / JPG).",
                     "• Interactive Threshold Slider (0.1 - 0.9).",
                     "• Visual 4-Panel Output & Laporan Severity."])

    add_clean_block(sl15, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.1),
                    "Arsitektur Cloudflare Tunnel",
                    ["• Daemon cloudflared outbound tunnel.",
                     "• Enkripsi SSL/TLS HTTPS otomatis.",
                     "• Custom URL: https://grad.mhisyam.com."])

    # =========================================================================
    # SLIDE 16: KESIMPULAN
    # =========================================================================
    sl16 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl16, "KESIMPULAN PENELITIAN", "Sintesis Capaian Penelitian Berdasarkan Tujuan & Pengujian")
    add_slide_footer(sl16, 16)

    add_clean_block(sl16, Inches(0.6), Inches(1.4), Inches(12.0), Inches(4.5),
                    "Kesimpulan Penelitian (Bab V — Penutup)",
                    ["1. Model Terbangun: Arsitektur U-Net + sCSE + EfficientNet-B3 berhasil dikembangkan. Pada data validasi diperoleh Precision 68,68%, Recall (Sensitivitas) 70,82%, dan Specificity 97,81%. Spesifisitas tinggi membuktikan model sangat jarang salah mengenali paru sehat sebagai pneumonia.",
                     "2. Aplikasi Web Gradio: Prototipe antarmuka web interaktif berhasil dibuat — menampilkan overlay segmentasi, peta panas Grad-CAM, analisis keparahan (5 tingkat), dan lateralitas infeksi dalam 1–3 detik.",
                     "3. Dual Lung Masking (PSPNet): Meminimalkan bias shortcut learning dari teks/penanda DICOM, meningkatkan kepercayaan klinis melalui Grad-CAM real-time.",
                     "4. Keterbatasan: Keterbatasan VRAM GPU (8 GB) mengharuskan gradient accumulation & AMP. Resolusi anotasi RSNA yang kasar (bounding box) membatasi presisi piksel. Belum ada validasi klinis langsung dengan dokter radiolog."])

    add_alert_block(sl16, Inches(0.6), Inches(6.1), Inches(12.0), Inches(0.65),
                    "Catatan Penting",
                    "Sistem ini berfungsi sebagai alat bantu skrining awal (triase), bukan pengganti keputusan akhir dokter spesialis radiologi.")

    # =========================================================================
    # SLIDE 17: SARAN PENGEMBANGAN
    # =========================================================================
    sl17 = prs.slides.add_slide(blank_layout)
    add_slide_header(sl17, "SARAN PENGEMBANGAN", "Rekomendasi untuk Peningkatan Performa & Penerapan Klinis Masa Depan")
    add_slide_footer(sl17, 17)

    add_clean_block(sl17, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.2),
                    "1. Rekomendasi Dataset & Model",
                    ["• Anotasi Piksel Asli: Menggunakan dataset rontgen dada dengan anotasi manual pixel-level langsung dari dokter spesialis radiologi, menggantikan konversi bounding box.",
                     "• Arsitektur Vision Transformer: Mengeksplorasi ViT-based segmentation (Swin-UNet, Medical SAM) untuk perbandingan kinerja."])

    add_clean_block(sl17, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.2),
                    "2. Rekomendasi Klinis & Sistem",
                    ["• Validasi Klinis Lokal: Melakukan uji coba lapangan menggunakan data pasien dari rumah sakit lokal Indonesia untuk menguji generalisasi model terhadap variasi mesin rontgen.",
                     "• Integrasi PACS: Mengintegrasikan aplikasi web dengan sistem Picture Archiving and Communication System (PACS) agar dapat digunakan langsung di lingkungan klinis nyata."],
                    title_color=C_PURPLE)

    # Save output
    prs.save(OUT_FILE)
    size_mb = os.path.getsize(OUT_FILE) / (1024 * 1024)
    print(f"✅  Native PPTX successfully generated: {OUT_FILE} ({size_mb:.2f} MB)")


if __name__ == "__main__":
    create_presentation()


